from flask import Flask, render_template, request, send_file, jsonify
from pptx import Presentation
import os
import zipfile

app = Flask(__name__)

UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "output"

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")

@app.route("/split_pptx", methods=["POST"])
def split_pptx():
    if "file" not in request.files:
        return jsonify({"error": "Nenhum arquivo enviado"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "Nenhum arquivo selecionado"}), 400

    filepath = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(filepath)

    # Carregar apresentação
    prs = Presentation(filepath)

    # Coletar intervalos de slides e nomes de arquivos
    slide_ranges = request.form.getlist("slideRanges[]")
    file_names = request.form.getlist("fileNames[]")
    slide_ranges = [tuple(map(int, r.split("-"))) for r in slide_ranges]

    zip_path = os.path.join(OUTPUT_FOLDER, "arquivos_divididos.zip")

    with zipfile.ZipFile(zip_path, "w") as zipf:
        for i, (start, end) in enumerate(slide_ranges):
            new_prs = Presentation()

            # Copiar slides no intervalo especificado
            for j in range(start - 1, end):
                slide_layout = new_prs.slide_layouts[0]  # Mantendo um layout básico
                new_slide = new_prs.slides.add_slide(slide_layout)

                for shape in prs.slides[j].shapes:
                    if hasattr(shape, "text"):
                        new_slide.shapes.title.text = shape.text

            output_filename = f"{file_names[i]}.pptx"
            output_path = os.path.join(OUTPUT_FOLDER, output_filename)
            new_prs.save(output_path)

            zipf.write(output_path, arcname=output_filename)

    return send_file(zip_path, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
